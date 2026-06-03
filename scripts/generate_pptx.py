#!/usr/bin/env python3
"""
Generate PPTX from slides/slide_contents.md
Requires: python-pptx
Usage: python3 scripts/generate_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
MD_PATH = os.path.join(ROOT, 'slides', 'slide_contents.md')
OUT_PPTX = os.path.join(ROOT, 'DMA_training_vs_STM32F4.pptx')


def load_slides(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        text = f.read()
    parts = [p.strip() for p in text.split('---') if p.strip()]
    slides = []
    for part in parts:
        lines = [l.rstrip() for l in part.splitlines() if l.strip()]
        if not lines:
            continue
        # First non-empty line may be Slide X: Title
        title_line = lines[0]
        # Extract title after colon if present
        if ':' in title_line:
            _, title = title_line.split(':', 1)
            title = title.strip()
        else:
            title = title_line
        # Collect bullet points (lines that start with '-' or not '备注:')
        bullets = []
        notes_lines = []
        in_notes = False
        for ln in lines[1:]:
            if ln.startswith('备注:'):
                in_notes = True
                notes_lines.append(ln[len('备注:'):].strip())
            elif in_notes:
                notes_lines.append(ln.strip())
            else:
                # remove leading '- '
                if ln.startswith('- '):
                    bullets.append(ln[2:].strip())
                else:
                    bullets.append(ln.strip())
        slides.append({'title': title, 'bullets': bullets, 'notes': '\n'.join(notes_lines)})
    return slides


def create_presentation(slides, out_path):
    prs = Presentation()
    # Set slide width/height defaults if needed
    # Title slide layout
    title_slide_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    first = True
    for s in slides:
        if first:
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = s['title'] if s['title'] else 'DMA Training'
            # Put first bullet as subtitle if exists
            subtitle.text = '\n'.join(s['bullets']) if s['bullets'] else ''
            # style
            title.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle.text_frame.paragraphs[0].font.size = Pt(14)
            first = False
        else:
            slide = prs.slides.add_slide(body_layout)
            title = slide.shapes.title
            title.text = s['title']
            title.text_frame.paragraphs[0].font.size = Pt(24)
            # body textbox
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, b in enumerate(s['bullets']):
                p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
                p.text = b
                p.level = 0
                p.font.size = Pt(18)
                p.font.name = 'Calibri'
                # dark blue color
                p.font.color.rgb = RGBColor(10, 49, 97)
            # add notes
        if s['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = s['notes']

    # Save
    prs.save(out_path)
    print('Saved PPTX to', out_path)


if __name__ == '__main__':
    if not os.path.exists(MD_PATH):
        print('Could not find', MD_PATH)
    else:
        slides = load_slides(MD_PATH)
        create_presentation(slides, OUT_PPTX)
